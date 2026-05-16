#!/usr/bin/env python
"""
答辩PPT自动生成器
基于答辩PPT大纲.md，用python-pptx从零生成15页毕业论文答辩PPT。
所有配图路径、文字内容、配色均可在此文件中集中修改，重新运行即可生成新版PPT。

用法：
    python gen_defense_ppt.py

输出：
    tmp/答辩_生成版.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image as PILImage

# ============================================================
# 路径配置 — 需要修改图片路径时改这里
# ============================================================
BASE = os.path.dirname(os.path.abspath(__file__))
DIAGRAMS = os.path.join(BASE, 'docs', 'diagrams')
FIGURES = os.path.join(BASE, 'MLP', 'outputs', 'training_jobs',
                       '20260420132754-6a846b-论文数据及图像', 'figures')
OUTPUT = os.path.join(BASE, 'tmp', '答辩_生成版.pptx')

# 图片文件映射
IMG = {
    'architecture': os.path.join(DIAGRAMS, '系统框架图.png'),
    'er_diagram': os.path.join(DIAGRAMS, '系统E-R图.png'),
    'claim_dist': os.path.join(DIAGRAMS, '数据理赔金额分布图.png'),
    'mlp_structure': os.path.join(DIAGRAMS, 'InsuranceMLP网络结构图.png'),
    'confusion_matrix': os.path.join(FIGURES, 'confusion_matrix.png'),
    'loss_curve': os.path.join(FIGURES, 'loss_curve.png'),
    'pr_auc_curve': os.path.join(FIGURES, 'val_pr_auc_curve.png'),
    'accuracy_curve': os.path.join(FIGURES, 'accuracy_curve.png'),
}

# ============================================================
# 配色方案
# ============================================================
C = {
    'dark_blue':   RGBColor(0x1A, 0x3C, 0x6E),
    'mid_blue':    RGBColor(0x40, 0x80, 0xC0),
    'light_blue':  RGBColor(0xD6, 0xE4, 0xF0),
    'pale_blue':   RGBColor(0xF0, 0xF4, 0xFF),
    'dark_gray':   RGBColor(0x33, 0x33, 0x33),
    'mid_gray':    RGBColor(0x66, 0x66, 0x66),
    'light_gray':  RGBColor(0xE0, 0xE0, 0xE0),
    'white':       RGBColor(0xFF, 0xFF, 0xFF),
    'black':       RGBColor(0x00, 0x00, 0x00),
    'accent_orange': RGBColor(0xE0, 0x7A, 0x2F),
    'accent_green':  RGBColor(0x2E, 0x7D, 0x32),
}

# ============================================================
# 尺寸常量 (16:9 宽屏: 20.0" x 11.25")
# ============================================================
SLIDE_W = Inches(20.0)
SLIDE_H = Inches(11.25)
MARGIN_L = Inches(1.0)
MARGIN_R = Inches(19.0)
CONTENT_W = Inches(18.0)
TITLE_Y = Inches(0.7)
SUBTITLE_Y = Inches(1.15)
PAGE_NUM_X = Inches(18.5)
PAGE_NUM_Y = Inches(0.7)

# ============================================================
# 工具函数
# ============================================================

def new_prs():
    """创建空白16:9演示文稿"""
    prs = Presentation()
    prs.slide_width = int(SLIDE_W)
    prs.slide_height = int(SLIDE_H)
    return prs


def add_blank_slide(prs):
    """添加空白幻灯片"""
    layout = prs.slide_layouts[6]  # blank layout
    return prs.slides.add_slide(layout)


def add_textbox(slide, left, top, width, height, text, font_size=Pt(14),
                color=C['dark_gray'], bold=False, alignment=PP_ALIGN.LEFT,
                font_name='Microsoft YaHei', anchor=MSO_ANCHOR.TOP):
    """添加文本框并返回文本框对象"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    txBox.text_frame.auto_size = None
    tf = txBox.text_frame
    tf.paragraphs[0].alignment = alignment
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    # Set East Asian font
    for run in p.runs:
        rPr = run._r.get_or_add_rPr()
        from lxml import etree
        ea = etree.SubElement(rPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}ea')
        ea.set('typeface', font_name)
    try:
        tf.paragraphs[0].font.name = font_name
    except Exception:
        pass
    return txBox


def add_multiline(slide, left, top, width, lines, font_size=Pt(12),
                  color=C['dark_gray'], line_spacing=Pt(20), bold_first=False):
    """添加多行文本。lines 是 [(text, is_bold), ...] 或 [str, ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, str):
            text, is_bold = line, False
        else:
            text, is_bold = line
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = font_size
        p.font.color.rgb = color
        p.font.bold = is_bold or (bold_first and i == 0)
        p.font.name = 'Microsoft YaHei'
        p.space_after = Pt(4)
    return txBox


def add_title_bar(slide, title, page_num, subtitle=None):
    """统一标题栏：深蓝横条 + 标题文字 + 页码"""
    # 深蓝顶条
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.25)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = C['dark_blue']
    bar.line.fill.background()

    # 标题
    add_textbox(slide, MARGIN_L, TITLE_Y, Inches(16.0), Inches(0.55),
                title, font_size=Pt(26), color=C['dark_blue'], bold=True)

    # 副标题
    if subtitle:
        add_textbox(slide, MARGIN_L, SUBTITLE_Y, Inches(17.0), Inches(0.4),
                    subtitle, font_size=Pt(14), color=C['mid_gray'])

    # 页码
    add_textbox(slide, PAGE_NUM_X, PAGE_NUM_Y, Inches(0.6), Inches(0.35),
                str(page_num), font_size=Pt(12), color=C['mid_gray'],
                alignment=PP_ALIGN.RIGHT)


def add_image_safe(slide, path_key_or_path, left, top, width=None, height=None):
    """智能添加图片。支持 IMG 字典的 key 或直接路径。width/height 指定一个即可保持比例。"""
    path = IMG.get(path_key_or_path, path_key_or_path)
    if not os.path.exists(path):
        print(f"  WARNING: Image not found: {path}")
        return None
    img = PILImage.open(path)
    iw, ih = img.size
    if width and height:
        return slide.shapes.add_picture(path, left, top, width, height)
    elif width:
        h = int(width * ih / iw)
        return slide.shapes.add_picture(path, left, top, width, h)
    elif height:
        w = int(height * iw / ih)
        return slide.shapes.add_picture(path, left, top, w, height)
    else:
        return slide.shapes.add_picture(path, left, top)


def add_card(slide, left, top, width, height, title, body_lines,
             title_color=C['dark_blue'], bg_color=C['pale_blue']):
    """添加带标题和内容的卡片"""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    card.line.color.rgb = C['light_gray']
    card.line.width = Pt(0.5)

    # 标题
    add_textbox(slide, left + Inches(0.3), top + Inches(0.15),
                width - Inches(0.6), Inches(0.4),
                title, font_size=Pt(16), color=title_color, bold=True)

    # 内容
    add_multiline(slide, left + Inches(0.3), top + Inches(0.55),
                  width - Inches(0.6), body_lines,
                  font_size=Pt(11), color=C['dark_gray'], line_spacing=Pt(18))


def add_section_divider(slide, section_num, section_title, left, top, width):
    """添加章节分隔条"""
    add_textbox(slide, left, top, Inches(0.6), Inches(0.4),
                f'0{section_num}', font_size=Pt(22), color=C['mid_blue'], bold=True)
    add_textbox(slide, left + Inches(0.7), top + Inches(0.05),
                width - Inches(0.7), Inches(0.35),
                section_title, font_size=Pt(18), color=C['dark_blue'], bold=True)
    # 下划线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left + Inches(0.7), top + Inches(0.45),
        width - Inches(0.7), Inches(0.015)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = C['mid_blue']
    line.line.fill.background()


# ============================================================
# 各页构建函数
# ============================================================

def build_cover(prs):
    """第1页：封面"""
    slide = add_blank_slide(prs)

    # 深蓝背景块（上半部分）
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = C['dark_blue']
    bg.line.fill.background()

    # 标题
    add_textbox(slide, Inches(2.0), Inches(3.0), Inches(16.0), Inches(1.5),
                '基于深度学习的\n车险理赔预测系统设计与实现',
                font_size=Pt(36), color=C['white'], bold=True, alignment=PP_ALIGN.CENTER)

    # 英文标题
    add_textbox(slide, Inches(2.0), Inches(5.0), Inches(16.0), Inches(0.5),
                'Design and Implementation of Auto Insurance Claim\nPrediction System Based on Deep Learning',
                font_size=Pt(14), color=RGBColor(0xBB, 0xCC, 0xDD), alignment=PP_ALIGN.CENTER)

    # 底部信息
    add_textbox(slide, Inches(2.0), Inches(8.5), Inches(16.0), Inches(0.35),
                '重庆邮电大学  ·  人工智能学院  ·  数据科学与大数据技术',
                font_size=Pt(16), color=C['dark_blue'], alignment=PP_ALIGN.CENTER)

    info_lines = [
        ('学生：钱信宇', False),
        ('指导教师：陈力 讲师', False),
        ('2026年5月', False),
    ]
    y = Inches(9.2)
    for text, bold in info_lines:
        add_textbox(slide, Inches(7.0), y, Inches(6.0), Inches(0.3),
                    text, font_size=Pt(14), color=C['mid_gray'], alignment=PP_ALIGN.CENTER)
        y += Inches(0.35)


def build_outline(prs):
    """第2页：目录"""
    slide = add_blank_slide(prs)
    add_title_bar(slide, '目录', 2, subtitle=None)

    sections = [
        ('01', '研究背景与目标', '问题定义、研究意义、核心贡献'),
        ('02', '系统架构与数据库设计', '三层架构、技术栈、E-R图'),
        ('03', '数据处理与特征工程', '缺失值处理、异常值截断、特征构建'),
        ('04', '模型构建与训练策略', '残差MLP、加权损失、自动阈值'),
        ('05', '实验结果与消融分析', '混淆矩阵、训练曲线、消融验证'),
        ('06', '系统演示与总结', '功能展示、量化成果、未来展望'),
    ]

    y_start = Inches(2.5)
    for i, (num, title, desc) in enumerate(sections):
        y = y_start + Inches(i * 1.3)
        # 序号
        add_textbox(slide, Inches(1.5), y, Inches(0.8), Inches(0.6),
                    num, font_size=Pt(32), color=C['mid_blue'], bold=True)
        # 标题
        add_textbox(slide, Inches(2.6), y, Inches(14.0), Inches(0.4),
                    title, font_size=Pt(20), color=C['dark_blue'], bold=True)
        # 描述
        add_textbox(slide, Inches(2.6), y + Inches(0.45), Inches(14.0), Inches(0.3),
                    desc, font_size=Pt(13), color=C['mid_gray'])
        # 分隔线
        if i < len(sections) - 1:
            sep = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(2.6), y + Inches(1.0),
                Inches(14.5), Inches(0.005)
            )
            sep.fill.solid()
            sep.fill.fore_color.rgb = C['light_gray']
            sep.line.fill.background()

    # 底部署名
    add_textbox(slide, MARGIN_L, Inches(10.2), Inches(18.0), Inches(0.3),
                '模型研究 + 系统工程', font_size=Pt(12), color=C['mid_gray'],
                alignment=PP_ALIGN.CENTER)


def build_background(prs):
    """第3页：研究背景与问题"""
    slide = add_blank_slide(prs)
    add_title_bar(slide, '研究背景与问题', 3, subtitle='车险行业正从经验驱动向数据驱动转型')

    # 左半栏：痛点
    add_section_divider(slide, 1, '行业痛点', Inches(1.0), Inches(2.0), Inches(8.5))
    pain_points = [
        '我国机动车保有量超4亿辆，车险保费规模破万亿，理赔案件量持续攀升',
        '传统人工审核依赖核保员经验，面对海量、多维、异构数据时效率与准确性不足',
        '理赔预测不准确 → 赔付成本上升 + 客户满意度下降 + 欺诈风险增加',
    ]
    y = Inches(2.7)
    for i, point in enumerate(pain_points):
        prefix = '▸ '
        add_textbox(slide, Inches(1.5), y, Inches(8.3), Inches(0.6),
                    prefix + point, font_size=Pt(13), color=C['dark_gray'])
        y += Inches(0.75)

    # 右半栏：机遇
    add_section_divider(slide, 2, '技术机遇', Inches(10.5), Inches(2.0), Inches(8.5))
    opportunities = [
        '深度学习能从多维结构化特征中自动学习非线性风险模式，识别潜在高风险保单',
        '将"模型训练→理赔预测→结果解释→统计反馈"整合至统一信息平台，降低算法应用门槛',
        '为保险企业提供一套可落地的数字化风控原型方案',
    ]
    y = Inches(2.7)
    for i, opp in enumerate(opportunities):
        prefix = '▸ '
        add_textbox(slide, Inches(11.0), y, Inches(8.3), Inches(0.6),
                    prefix + opp, font_size=Pt(13), color=C['dark_gray'])
        y += Inches(0.75)

    # 底部关键数字
    nums = [('105,555', '条保险记录'), ('30', '个原始特征'), ('18.6%', '理赔占比')]
    x_start = Inches(3.0)
    for i, (num, label) in enumerate(nums):
        x = x_start + Inches(i * 5.0)
        add_textbox(slide, x, Inches(9.2), Inches(4.0), Inches(0.6),
                    num, font_size=Pt(30), color=C['dark_blue'], bold=True,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, Inches(9.8), Inches(4.0), Inches(0.3),
                    label, font_size=Pt(12), color=C['mid_gray'], alignment=PP_ALIGN.CENTER)


def build_contributions(prs):
    """第4页：研究目标与核心贡献"""
    slide = add_blank_slide(prs)
    add_title_bar(slide, '研究目标与核心贡献', 4,
                  subtitle='目标：设计并实现一个结合深度学习预测模型的车险理赔管理平台')

    contributions = [
        ('残差式 MLP 预测模型',
         ['设计含残差连接的多层感知机网络', '配合 LayerNorm + GELU + Dropout', '解决深层网络退化与过拟合问题']),
        ('类别不平衡训练策略',
         ['加权 BCE 损失 + Fβ 自动阈值搜索', 'Recall ≥ 0.830 业务硬约束', '将模型输出概率与业务决策边界解耦']),
        ('前后端分离业务平台',
         ['Spring Boot + Vue3 + FastAPI 三端协同', '数据管理、模型训练、理赔预测、统计分析', '形成完整的业务闭环']),
        ('系统化消融验证',
         ['加权损失 vs 无加权 × 自动阈值 vs 固定阈值', '网络深度 2→5 层的三维对比实验', '验证每个设计选择的必要性']),
    ]

    card_w = Inches(4.2)
    card_h = Inches(3.8)
    gap = Inches(0.3)
    x_start = Inches(0.9)
    y = Inches(2.2)

    for i, (title, lines) in enumerate(contributions):
        x = x_start + i * (card_w + gap)
        add_card(slide, x, y, card_w, card_h, title, lines,
                 title_color=C['dark_blue'], bg_color=C['pale_blue'])


def build_architecture(prs):
    """第5页：系统总体架构"""
    slide = add_blank_slide(prs)
    add_title_bar(slide, '系统总体架构', 5,
                  subtitle='三层架构：表示层 → 业务层 → 数据层，预测服务通过 API 解耦')

    # 左侧：架构图
    if os.path.exists(IMG['architecture']):
        add_image_safe(slide, 'architecture',
                       Inches(0.5), Inches(2.2), height=Inches(7.8))

    # 右侧：各层说明
    layers = [
        ('表示层', 'Vue3 + Element Plus + ECharts\nAxios 统一封装 HTTP 请求\n响应式数据绑定与路由管理', C['mid_blue']),
        ('业务层', 'Spring Boot + MyBatis\nREST API + 统一 Result 响应\n会话认证 + 角色权限控制', C['dark_blue']),
        ('数据层', 'MySQL 关系型数据库\n用户、保单、理赔、预测等核心表\n定期备份 + 加密传输', C['accent_green']),
        ('预测服务', 'FastAPI + PyTorch\n封装训练好的残差 MLP 模型\nJSON 接口：概率 / 标签 / 风险等级', C['accent_orange']),
    ]

    y = Inches(2.2)
    for title, desc, color in layers:
        # 色块标签
        tag = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.5), y, Inches(1.5), Inches(0.35)
        )
        tag.fill.solid()
        tag.fill.fore_color.rgb = color
        tag.line.fill.background()
        tf = tag.text_frame
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.color.rgb = C['white']
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        add_textbox(slide, Inches(13.2), y - Inches(0.05), Inches(5.8), Inches(1.0),
                    desc, font_size=Pt(10), color=C['dark_gray'])
        y += Inches(1.65)


def build_database(prs):
    """第6页：数据库设计"""
    slide = add_blank_slide(prs)
    add_title_bar(slide, '数据库设计', 6, subtitle='核心实体关系与数据表结构')

    # E-R图
    if os.path.exists(IMG['er_diagram']):
        add_image_safe(slide, 'er_diagram',
                       Inches(0.5), Inches(2.2), width=Inches(11.0))

    # 右侧说明
    entities = [
        ('sys_user（用户表）', '用户ID、角色、联系方式、注册时间'),
        ('motor_insurance（保单表）', '保单号、投保人/车辆信息、保费、风险类型'),
        ('claim_record（理赔记录）', '理赔ID、关联保单、事故描述、理赔金额'),
        ('vehicle_info（车辆信息）', '车牌号、车型、车龄、燃油类型'),
        ('insur_pred（预测记录）', '关联保单、预测概率、风险等级、模型版本'),
    ]
    y = Inches(2.2)
    for title, desc in entities:
        add_textbox(slide, Inches(12.0), y, Inches(7.0), Inches(0.3),
                    title, font_size=Pt(12), color=C['dark_blue'], bold=True)
        add_textbox(slide, Inches(12.0), y + Inches(0.3), Inches(7.0), Inches(0.3),
                    desc, font_size=Pt(10), color=C['mid_gray'])
        y += Inches(0.75)

    # 底部注释
    add_textbox(slide, Inches(12.0), Inches(8.5), Inches(7.0), Inches(0.8),
                'train_data 表由保单、理赔、车辆三表\n按 ID 关联聚合生成，供模型训练使用',
                font_size=Pt(10), color=C['mid_gray'])


def build_preprocessing(prs):
    """第7页：数据预处理"""
    slide = add_blank_slide(prs)
    add_title_bar(slide, '数据预处理', 7,
                  subtitle='数据集：105,555 条记录 × 30 个特征，理赔占比 18.6%（正负比 ≈ 1:4.37）')

    # 四步流水线
    steps = [
        ('① 缺失值处理', 'Data_Lapse（66.7%缺失）\n→ 直接删除字段', 'Length / Type_fuel\n→ 按车型分组众数填充'),
        ('② 异常值截断', '保费、车价、功率等偏态字段\n→ 99.9 分位数截断', '保留合理高值\n剔除极端离群点'),
        ('③ 不平衡处理', '分层采样 70/15/15\n划分训练/验证/测试', '训练时使用\npos_weight 类别权重'),
        ('④ 编码转换', '分类变量 → 整数编码', '保持特征含义\n适配 MLP 输入'),
    ]

    x_start = Inches(0.7)
    card_w = Inches(4.2)
    card_h = Inches(3.0)
    gap = Inches(0.25)

    for i, (title, line1, line2) in enumerate(steps):
        x = x_start + i * (card_w + gap)
        # 卡片背景
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.3), card_w, card_h
        )
        card.fill.solid()
        card.fill.fore_color.rgb = C['pale_blue']
        card.line.color.rgb = C['light_gray']
        card.line.width = Pt(0.5)

        # 步骤标题
        add_textbox(slide, x + Inches(0.2), Inches(2.5), card_w - Inches(0.4), Inches(0.4),
                    title, font_size=Pt(15), color=C['dark_blue'], bold=True)
        # 内容
        add_textbox(slide, x + Inches(0.2), Inches(3.1), card_w - Inches(0.4), Inches(1.0),
                    line1, font_size=Pt(11), color=C['dark_gray'])
        add_textbox(slide, x + Inches(0.2), Inches(4.1), card_w - Inches(0.4), Inches(0.8),
                    line2, font_size=Pt(11), color=C['dark_gray'])

    # 底部：分布图
    if os.path.exists(IMG['claim_dist']):
        add_image_safe(slide, 'claim_dist',
                       Inches(2.5), Inches(5.6), height=Inches(3.8))

    # 底部说明
    add_textbox(slide, MARGIN_L, Inches(10.5), Inches(18.0), Inches(0.3),
                '预处理目的：保证输入数据一致性，控制极端值干扰，缓解类别不平衡偏置',
                font_size=Pt(11), color=C['mid_gray'], alignment=PP_ALIGN.CENTER)


def build_features(prs):
    """第8页：特征工程"""
    slide = add_blank_slide(prs)
    add_title_bar(slide, '特征工程', 8,
                  subtitle='将业务时间、物理属性和历史合同信息统一为可训练的数值特征矩阵')

    # 三列特征卡片
    categories = [
        ('时间锚点特征', C['dark_blue'],
         ['以合同起始日期为统一参考点',
          'Date_start_contract_days',
          'insured_age_years',
          'vehicle_age_years',
          '→ 将绝对日期转为相对天数']),
        ('连续数值特征', C['accent_green'],
         ['经 99.9 分位数截断后保留',
          'Power（功率）',
          'Weight（车重）',
          'Value_vehicle（车辆价值）',
          'Premium（保费）']),
        ('历史与合同特征', C['accent_orange'],
         ['描述保单的静态属性与历史状态',
          'Type_risk（风险类型）',
          'Payment（付款方式）',
          'Type_fuel（燃油类型）',
          '→ 整数编码后输入模型']),
    ]

    col_w = Inches(5.8)
    gap = Inches(0.4)
    x_start = Inches(0.8)

    for i, (title, color, items) in enumerate(categories):
        x = x_start + i * (col_w + gap)
        # 标题卡片
        header = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.2), col_w, Inches(0.5)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = color
        header.line.fill.background()
        tf = header.text_frame
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = C['white']
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 特征内容
        add_multiline(slide, x + Inches(0.3), Inches(3.0), col_w - Inches(0.6),
                      items, font_size=Pt(12), color=C['dark_gray'])

    # 底部总结
    add_textbox(slide, MARGIN_L, Inches(9.5), Inches(18.0), Inches(0.6),
                '特征工程核心思路：不是简单罗列原始字段，而是按业务逻辑分组构建，'
                '将时间、物理属性、历史合同信息统一为结构化特征矩阵',
                font_size=Pt(12), color=C['mid_gray'], alignment=PP_ALIGN.CENTER)


def build_model_structure(prs):
    """第9页：模型结构设计"""
    slide = add_blank_slide(prs)
    add_title_bar(slide, '模型结构设计', 9,
                  subtitle='残差式 MLP（Multi-Layer Perceptron with Residual Connections）')

    # MLP结构图
    if os.path.exists(IMG['mlp_structure']):
        add_image_safe(slide, 'mlp_structure',
                       Inches(0.5), Inches(2.3), width=Inches(12.5))

    # 右侧设计要点
    design_points = [
        ('残差连接', '每个 ResidualBlock 包含\nFC→LayerNorm→GELU→Dropout→FC→LayerNorm\n与捷径支路相加后经 GELU 激活'),
        ('捷径对齐', '输入输出维度不一致时\n通过无偏置线性投影对齐\nnn.Linear(in, out, bias=False)'),
        ('输出层', 'Sigmoid 将 logit 映射为 [0,1]\n表示"未来一年内发生理赔"的概率'),
        ('网络结构', 'input → 256 → 512 → 512 → 256 → 128 → 1\n共 5 个残差块，总参数量适中'),
    ]

    y = Inches(2.3)
    for title, desc in design_points:
        add_textbox(slide, Inches(13.5), y, Inches(5.8), Inches(0.3),
                    title, font_size=Pt(14), color=C['dark_blue'], bold=True)
        add_textbox(slide, Inches(13.5), y + Inches(0.35), Inches(5.8), Inches(1.2),
                    desc, font_size=Pt(10), color=C['dark_gray'])
        y += Inches(1.65)


def build_training_strategy(prs):
    """第10页：训练策略"""
    slide = add_blank_slide(prs)
    add_title_bar(slide, '训练策略', 10,
                  subtitle='针对类别不平衡场景的四项核心策略')

    cards = [
        ('加权 BCE 损失',
         ['使用 pos_weight 加权的 BCEWithLogitsLoss',
          '增大正样本（理赔）的误分类惩罚',
          '让模型"更关注少数类的错误"',
          'pos_weight ≈ 4.37（正负样本反比）']),
        ('Fβ 自动阈值搜索',
         ['在 [0.05, 0.95] 区间扫描 200 个候选阈值',
          '优化目标：Fβ（β=1.3，偏向召回率）',
          '硬约束：Recall ≥ 0.830',
          '作用：模型输出概率与业务决策边界解耦']),
        ('AdamW + Cosine Warmup',
         ['AdamW 解耦权重衰减与自适应学习率',
          '比标准 Adam 泛化能力更强',
          'Cosine Warmup 从低 LR 逐步升温再衰减',
          '训练初期更稳定，收敛路径更平滑']),
        ('AMP + 梯度裁剪 + Early Stopping',
         ['AMP 混合精度加速训练、节省显存',
          '梯度裁剪防止梯度爆炸',
          '以验证集 AUC 为早停指标',
          'AUC 对阈值不敏感，比 F1 更稳定']),
    ]

    card_w = Inches(4.3)
    card_h = Inches(4.0)
    gap = Inches(0.2)
    x_start = Inches(0.7)

    for i, (title, lines) in enumerate(cards):
        x = x_start + i * (card_w + gap)
        add_card(slide, x, Inches(2.2), card_w, card_h, title, lines,
                 title_color=C['dark_blue'])

    # 底部业务含义
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(6.6), Inches(18.6), Inches(0.6)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = C['dark_blue']
    bar.line.fill.background()
    tf = bar.text_frame
    tf.paragraphs[0].text = ('业务含义：模型只管输出概率，阈值负责划定决策边界。'
                             '在控制理赔漏检率（Recall ≥ 0.83）的前提下适度接受误报'
                             '——契合保险风险管理的业务逻辑')
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = C['white']
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER


def build_results(prs):
    """第11页：实验结果"""
    slide = add_blank_slide(prs)
    add_title_bar(slide, '实验结果', 11,
                  subtitle='测试集 10,557 条 | 模型在保证多数类精度的前提下，对理赔类实现 83% 召回率')

    # 左侧三大指标
    metrics = [
        ('0.9604', 'AUC', '模型排序能力优异'),
        ('0.8305', 'Claim Recall', '正样本（理赔）召回率'),
        ('0.9121', 'Weighted F1', '加权平均 F1-score'),
    ]
    y = Inches(2.3)
    for value, label, desc in metrics:
        add_textbox(slide, Inches(0.9), y, Inches(5.5), Inches(0.7),
                    value, font_size=Pt(36), color=C['dark_blue'], bold=True)
        add_textbox(slide, Inches(0.9), y + Inches(0.7), Inches(5.5), Inches(0.3),
                    label, font_size=Pt(14), color=C['mid_gray'], bold=True)
        add_textbox(slide, Inches(0.9), y + Inches(1.0), Inches(5.5), Inches(0.25),
                    desc, font_size=Pt(10), color=C['mid_gray'])
        y += Inches(1.6)

    # 右侧分类报告表格区域
    # 用文本框模拟表格
    table_data = [
        ['类别', 'Precision', 'Recall', 'F1-score', '样本数'],
        ['No Claim', '0.9599', '0.9280', '0.9437', '8,592'],
        ['Claim', '0.7250', '0.8305', '0.7742', '1,965'],
        ['加权平均', '0.9162', '0.9098', '0.9121', '10,557'],
    ]

    table_left = Inches(7.0)
    table_top = Inches(2.3)
    col_widths = [Inches(1.8), Inches(1.8), Inches(1.8), Inches(2.0), Inches(1.6)]
    row_height = Inches(0.45)

    for r, row in enumerate(table_data):
        x_pos = table_left
        for c, cell in enumerate(row):
            is_header = (r == 0)
            is_last_row = (r == len(table_data) - 1)
            bg_color = C['dark_blue'] if is_header else (C['pale_blue'] if is_last_row else C['white'])
            text_color = C['white'] if is_header else C['dark_gray']

            cell_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                x_pos, table_top + r * row_height,
                col_widths[c], row_height
            )
            cell_shape.fill.solid()
            cell_shape.fill.fore_color.rgb = bg_color
            cell_shape.line.color.rgb = C['light_gray']
            cell_shape.line.width = Pt(0.5)

            tf = cell_shape.text_frame
            tf.paragraphs[0].text = cell
            tf.paragraphs[0].font.size = Pt(11) if not is_header else Pt(12)
            tf.paragraphs[0].font.color.rgb = text_color
            tf.paragraphs[0].font.bold = is_header or is_last_row
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            x_pos += col_widths[c]

    # 混淆矩阵图
    if os.path.exists(IMG['confusion_matrix']):
        add_image_safe(slide, 'confusion_matrix',
                       Inches(7.0), Inches(4.6), height=Inches(3.2))

    # 底部总结
    add_textbox(slide, MARGIN_L, Inches(10.5), Inches(18.0), Inches(0.3),
                '关键解读：No Claim 类精确率与召回率均 > 92%；Claim 类召回率 83% → 约 17% 高风险保单被漏判',
                font_size=Pt(11), color=C['mid_gray'], alignment=PP_ALIGN.CENTER)


def build_training_curves(prs):
    """第12页：训练过程与消融分析"""
    slide = add_blank_slide(prs)
    add_title_bar(slide, '训练过程与消融分析', 12,
                  subtitle='训练收敛验证 + 三维度消融实验')

    # 左侧：训练曲线
    if os.path.exists(IMG['loss_curve']):
        add_image_safe(slide, 'loss_curve',
                       Inches(0.5), Inches(2.2), height=Inches(3.5))
    if os.path.exists(IMG['pr_auc_curve']):
        add_image_safe(slide, 'pr_auc_curve',
                       Inches(0.5), Inches(6.0), height=Inches(3.5))

    # 右侧：消融分析
    add_section_divider(slide, 1, '消融实验结论', Inches(10.5), Inches(2.2), Inches(8.8))

    ablation_items = [
        ('损失函数',
         '加权 vs 无加权',
         'F1: 0.7046 → 0.7550\n召回率: 0.8198'),
        ('阈值策略',
         '固定 0.5 vs 自动搜索',
         '固定阈值下加权损失反而降低 F1\n说明自动阈值是必要配套'),
        ('网络深度',
         '2 层 / 3 层 / 4 层 / 5 层',
         'AUC: 0.9577 → 0.9589\n边际增益递减，5 层最优'),
    ]

    y = Inches(2.9)
    for title, compare, result in ablation_items:
        # 标题
        add_textbox(slide, Inches(10.8), y, Inches(8.5), Inches(0.3),
                    title, font_size=Pt(14), color=C['dark_blue'], bold=True)
        # 对比
        add_textbox(slide, Inches(10.8), y + Inches(0.35), Inches(4.0), Inches(0.25),
                    compare, font_size=Pt(10), color=C['mid_gray'])
        # 结果
        add_textbox(slide, Inches(15.0), y + Inches(0.3), Inches(4.3), Inches(0.5),
                    result, font_size=Pt(10), color=C['accent_green'])
        y += Inches(1.15)

    # 最终方案
    add_textbox(slide, Inches(10.5), Inches(8.5), Inches(8.8), Inches(1.0),
                '最终方案：5层残差结构(256→512→512→256→128)\n'
                '+ 加权BCE + Fβ(β=1.3)自动阈值 + Recall≥0.830\n'
                '+ AdamW + Cosine Warmup + Early Stopping(AUC)',
                font_size=Pt(11), color=C['dark_blue'], bold=True)


def build_system_demo_1(prs):
    """第13页：系统演示（一）"""
    slide = add_blank_slide(prs)
    add_title_bar(slide, '系统演示（一）：数据管理与模型训练', 13)

    # 功能要点
    functions = [
        ('业务数据管理',
         ['保险单信息增删改查 + 分页浏览',
          '理赔记录录入与关联查询',
          '车辆信息统一管理',
          '普通用户仅操作本人数据，管理员管理全局']),
        ('权限控制',
         ['ADMIN：管理全部数据 + 训练模型',
          'STUDENT（操作员）：管理本人数据 + 使用已有模型预测']),
        ('模型训练',
         ['管理员配置超参数（学习率 / 隐藏层 / 阈值策略）',
          '训练完成后查看指标和曲线',
          '支持多版本模型管理']),
    ]

    x_start = Inches(0.7)
    col_w = Inches(5.8)
    gap = Inches(0.3)
    for i, (title, lines) in enumerate(functions):
        x = x_start + i * (col_w + gap)
        add_card(slide, x, Inches(2.0), col_w, Inches(4.5), title, lines,
                 title_color=C['dark_blue'])

    # 底部说明
    add_textbox(slide, MARGIN_L, Inches(9.8), Inches(18.0), Inches(0.3),
                '注：建议替换为实际系统运行截图。可在系统中操作对应功能后截图放入此页和下页。',
                font_size=Pt(11), color=C['accent_orange'], alignment=PP_ALIGN.CENTER)


def build_system_demo_2(prs):
    """第14页：系统演示（二）"""
    slide = add_blank_slide(prs)
    add_title_bar(slide, '系统演示（二）：理赔预测与统计分析', 14)

    # 功能要点
    functions = [
        ('理赔预测',
         ['选择已有保单或手动填写信息',
          '系统返回：理赔概率 / 预测标签 / 风险等级',
          '提供局部特征解释：展示影响本次预测的主要因素',
          '预测结果持久化存储，支持历史查询']),
        ('统计分析',
         ['理赔预测结果的统计汇总',
          '风险等级分布可视化（ECharts）',
          '历史预测记录查询与筛选',
          '为业务决策提供数据支撑']),
    ]

    col_w = Inches(8.5)
    for i, (title, lines) in enumerate(functions):
        x = Inches(0.7) + i * (col_w + Inches(0.5))
        add_card(slide, x, Inches(2.0), col_w, Inches(4.5), title, lines)

    # 底部
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(8.5), Inches(18.6), Inches(0.6)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = C['dark_blue']
    bar.line.fill.background()
    tf = bar.text_frame
    tf.paragraphs[0].text = '业务闭环：数据录入 → 模型推理 → 结果解释 → 统计反馈'
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = C['white']
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 底部说明
    add_textbox(slide, MARGIN_L, Inches(9.8), Inches(18.0), Inches(0.3),
                '注：建议替换为实际系统运行截图（预测结果页 + 统计分析页各一张）。',
                font_size=Pt(11), color=C['accent_orange'], alignment=PP_ALIGN.CENTER)


def build_summary(prs):
    """第15页：总结与展望"""
    slide = add_blank_slide(prs)
    add_title_bar(slide, '总结与展望', 15)

    # 三栏布局
    columns = [
        ('主要成果', C['accent_green'],
         ['1. 构建残差式 MLP 理赔预测模型\n   AUC=0.9604，Recall=83.05%',
          '2. 提出加权损失+Fβ自动阈值\n   +Recall约束的训练方案',
          '3. 搭建前后端分离平台\n   实现数据管理到统计分析的闭环']),
        ('现有不足', C['accent_orange'],
         ['1. 仅做二分类（是否理赔）\n   未对理赔金额做回归估计',
          '2. 未与 XGBoost 等\n   主流基线模型做量化对比']),
        ('未来方向', C['dark_blue'],
         ['1. 将二分类升级为\n   赔付金额区间预测+多任务建模',
          '2. 接入车联网数据\n   （驾驶行为、轨迹等）丰富特征维度']),
    ]

    col_w = Inches(5.8)
    gap = Inches(0.4)
    for i, (title, color, lines) in enumerate(columns):
        x = Inches(0.8) + i * (col_w + gap)
        # 标题栏
        header = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.2), col_w, Inches(0.55)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = color
        header.line.fill.background()
        tf = header.text_frame
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.color.rgb = C['white']
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 内容卡片
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.9), col_w, Inches(5.5)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = C['pale_blue']
        card.line.color.rgb = C['light_gray']
        card.line.width = Pt(0.5)

        add_multiline(slide, x + Inches(0.3), Inches(3.3), col_w - Inches(0.6),
                      lines, font_size=Pt(12), color=C['dark_gray'])

    # 底部致谢
    add_textbox(slide, MARGIN_L, Inches(9.8), Inches(18.0), Inches(0.5),
                '感谢各位老师批评指正！',
                font_size=Pt(24), color=C['dark_blue'], bold=True,
                alignment=PP_ALIGN.CENTER)


# ============================================================
# 主函数
# ============================================================

def main():
    print('=' * 60)
    print('答辩PPT自动生成器')
    print('=' * 60)

    # 检查图片
    missing = [k for k, v in IMG.items() if not os.path.exists(v)]
    if missing:
        print(f'WARNING: 以下图片未找到: {missing}')
        print('PPT 将跳过缺失的图片继续生成')
    else:
        print('所有图片资源就绪 [OK]')

    prs = new_prs()
    print(f'幻灯片尺寸: {SLIDE_W/914400:.1f} x {SLIDE_H/914400:.1f} 英寸 (16:9)')
    print()

    builders = [
        ('封面', build_cover),
        ('目录', build_outline),
        ('研究背景与问题', build_background),
        ('研究目标与核心贡献', build_contributions),
        ('系统总体架构', build_architecture),
        ('数据库设计', build_database),
        ('数据预处理', build_preprocessing),
        ('特征工程', build_features),
        ('模型结构设计', build_model_structure),
        ('训练策略', build_training_strategy),
        ('实验结果', build_results),
        ('训练过程与消融分析', build_training_curves),
        ('系统演示（一）', build_system_demo_1),
        ('系统演示（二）', build_system_demo_2),
        ('总结与展望', build_summary),
    ]

    for name, builder in builders:
        builder(prs)
        print(f'  [{len(prs.slides):2d}/15] {name}')

    os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
    prs.save(OUTPUT)

    print()
    print(f'= 生成完成: {OUTPUT}')
    print(f'  共 {len(prs.slides)} 页幻灯片')
    print()
    print('提示：')
    print('  - 第13、14页系统演示需替换为实际截图')
    print('  - 配色/字体/内容可在脚本开头的配置区集中修改')
    print('  - 修改后重新运行 python gen_defense_ppt.py 即可')


if __name__ == '__main__':
    main()
